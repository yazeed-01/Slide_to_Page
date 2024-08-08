import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from docx import Document
from docx.enum.text import WD_COLOR_INDEX
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

def add_horizontal_line(doc):
    p = doc.add_paragraph()
    p.add_run().add_break()
    p.add_run().add_break()

    p = doc.add_paragraph()
    p.alignment = 1  

    run = p.add_run()
    border = OxmlElement('w:pBdr')
    top = OxmlElement('w:top')
    top.set(qn('w:val'), 'single')
    top.set(qn('w:sz'), '24')  
    top.set(qn('w:space'), '1')
    top.set(qn('w:color'), 'auto')
    border.append(top)
    p._element.get_or_add_pPr().append(border)

    p.add_run().add_break()
    p.add_run().add_break()

def convert_ppt_to_word(ppt_file, word_file):
    ppt = Presentation(ppt_file)
    doc = Document()

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            if shape == slide.shapes.title:
                title = shape.text_frame.text.strip()
                if title:
                    p = doc.add_paragraph()
                    run = p.add_run("Title: " + title)
                    run.bold = True
                    run.font.highlight_color = WD_COLOR_INDEX.TURQUOISE
            else:
                for idx, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    content = paragraph.text.strip()
                    if content:
                        p = doc.add_paragraph()
                        run = p.add_run(f"{idx}. {content}")

        add_horizontal_line(doc)

    doc.save(word_file)

def select_file():
    ppt_file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if ppt_file:
        save_file = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word files", "*.docx")])
        if save_file:
            try:
                convert_ppt_to_word(ppt_file, save_file)
                messagebox.showinfo("Success", "Conversion completed successfully!")
            except Exception as e:
                messagebox.showerror("Error", f"An error occurred: {e}")

# Create the main window
root = tk.Tk()
root.title("PPT to Word Converter")
root.geometry("400x200")  # Width x Height in pixels


select_button = tk.Button(root, text="Select PowerPoint File", command=select_file)
select_button.pack(pady=20)

root.mainloop()